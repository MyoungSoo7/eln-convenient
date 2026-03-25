#!/usr/bin/env python3
"""SI 산출물 PPTX 생성 (4건): 제안서, 주간보고서, 설계단계 완료보고서, 교육자료"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, datetime

OUT = os.path.dirname(os.path.abspath(__file__))
TODAY = datetime.date.today().isoformat()
C1=RGBColor(0x1E,0x40,0xAF); C2=RGBColor(0x3B,0x82,0xF6); CW=RGBColor(0xFF,0xFF,0xFF)
CD=RGBColor(0x1E,0x29,0x3B); CG=RGBColor(0x64,0x74,0x8B); CS=RGBColor(0x10,0xB9,0x81)

def mkprs():
    p=Presentation(); p.slide_width=Inches(13.333); p.slide_height=Inches(7.5); return p

def rect(s,l,t,w,h,fc=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    sh.line.fill.background(); return sh

def tb(s,l,t,w,h,txt,sz=14,c=CD,b=False,al=PP_ALIGN.LEFT):
    tx=s.shapes.add_textbox(l,t,w,h); tf=tx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.font.name='맑은 고딕'; p.alignment=al
    return tx

def hdr(s,txt,sub=""):
    rect(s,0,0,Inches(13.333),Inches(0.85),C1)
    tb(s,Inches(0.7),Inches(0.12),Inches(11),Inches(0.55),txt,24,CW,True)
    if sub: tb(s,Inches(0.7),Inches(0.55),Inches(11),Inches(0.3),sub,12,RGBColor(0xBF,0xDB,0xFE))

def bul(s,l,t,w,h,items,sz=12):
    tx=s.shapes.add_textbox(l,t,w,h); tf=tx.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=f"\u2022 {item}"; p.font.size=Pt(sz); p.font.color.rgb=CD; p.font.name='맑은 고딕'; p.space_after=Pt(4)

def tbl(s,l,t,data,cw):
    rows,cols=len(data),len(data[0])
    ts=s.shapes.add_table(rows,cols,l,t,Inches(sum(cw)),Inches(rows*0.35))
    tb2=ts.table
    for ci,w in enumerate(cw): tb2.columns[ci].width=Inches(w)
    for ri,row in enumerate(data):
        for ci,v in enumerate(row):
            c=tb2.cell(ri,ci); c.text=str(v)
            for p in c.text_frame.paragraphs:
                p.font.size=Pt(10); p.font.name='맑은 고딕'
                if ri==0: p.font.bold=True; p.font.color.rgb=CW
            if ri==0: c.fill.solid(); c.fill.fore_color.rgb=C1
            elif ri%2==0: c.fill.solid(); c.fill.fore_color.rgb=RGBColor(0xF1,0xF5,0xF9)

def cover(prs,title,sub):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    r1=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(13.333),Inches(7.5))
    r1.fill.gradient(); r1.fill.gradient_stops[0].color.rgb=C1; r1.fill.gradient_stops[1].color.rgb=RGBColor(0x0E,0x2A,0x6B); r1.line.fill.background()
    tb(s,Inches(1.5),Inches(2.2),Inches(10),Inches(1),title,44,CW,True,PP_ALIGN.CENTER)
    tb(s,Inches(1.5),Inches(3.3),Inches(10),Inches(0.5),sub,22,RGBColor(0x93,0xC5,0xFD),False,PP_ALIGN.CENTER)
    tb(s,Inches(1.5),Inches(4.5),Inches(10),Inches(0.4),"LabNote ELN 프로젝트",16,CW,False,PP_ALIGN.CENTER)
    tb(s,Inches(1.5),Inches(5.2),Inches(10),Inches(0.3),TODAY,12,RGBColor(0x93,0xC5,0xFD),False,PP_ALIGN.CENTER)

def sv(prs,name):
    prs.save(os.path.join(OUT,name)); print(f"  {name} ({len(prs.slides)} slides)")

# ════════════════════════════════════════════
# 02. 제안서
# ════════════════════════════════════════════
def g02():
    prs=mkprs(); cover(prs,"LabNote ELN 제안서","전자연구노트 협업 플랫폼 구축 제안")

    # 목차
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"목차")
    items=["1. 프로젝트 개요","2. 제안 배경 및 목표","3. 시스템 구성","4. 핵심 기능","5. 기술 아키텍처","6. 추진 일정","7. 투입 인력","8. 기대 효과"]
    bul(s,Inches(1),Inches(1.2),Inches(10),Inches(5),items,18)

    # 프로젝트 개요
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"1. 프로젝트 개요")
    tb(s,Inches(0.7),Inches(1.2),Inches(11),Inches(0.5),"연구실 실험 기록의 디지털 전환 — 전자서명 규정 준수, 협업 효율화, 통합 관리",14,CG)
    tbl(s,Inches(0.7),Inches(1.8),[["항목","내용"],["프로젝트명","LabNote ELN (전자연구노트 협업 플랫폼)"],["아키텍처","MSA 기반 9개 서비스, Docker Compose 온프레미스"],["기술 스택","React+Vite / Fastify / PostgreSQL / Redis / MinIO / OpenSearch"],["주요 기능","연구노트, 전자서명, 인벤토리, 예약, 검색, 내보내기, 감사로그"],["기간","3개월"]],[3,9])

    # 시스템 구성
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"3. 시스템 구성")
    svcs=[("api-gateway\n:8000","JWT, 프록시"),("auth\n:8001","인증, RBAC"),("eln\n:8002","연구노트"),("sig-audit\n:8003","서명, 감사"),("inventory\n:8004","시약/장비"),("scheduler\n:8005","예약"),("search\n:8006","통합검색"),("file\n:8008","파일"),("collab\n:8009","실시간 협업")]
    for i,(n,d) in enumerate(svcs):
        x=Inches(0.3)+Inches(i*1.43); y=Inches(1.5)
        rect(s,x,y,Inches(1.3),Inches(1.4),RGBColor(0xEF,0xF6,0xFF))
        tb(s,x+Inches(0.05),y+Inches(0.1),Inches(1.2),Inches(0.5),n,10,C1,True,PP_ALIGN.CENTER)
        tb(s,x+Inches(0.05),y+Inches(0.7),Inches(1.2),Inches(0.5),d,9,CG,False,PP_ALIGN.CENTER)

    # 핵심 기능
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"4. 핵심 기능")
    feats=[("연구노트","작성/편집/버전관리\n실시간 다중 사용자 협업"),("전자서명","SHA-256 해시체인\n무결성 검증, 규정 준수"),("인벤토리","시약/장비 재고 추적\n바코드, 만료/부족 알림"),("예약","장비/회의실 예약\n승인 워크플로, 충돌 방지"),("통합 검색","한국어 형태소 분석\n노트/템플릿/인벤토리"),("내보내기","PDF/ZIP 자동 생성\nPuppeteer, 비동기")]
    for i,(t,d) in enumerate(feats):
        c,r=i%3,i//3; x=Inches(0.5)+Inches(c*4.2); y=Inches(1.2)+Inches(r*2.8)
        rect(s,x,y,Inches(3.9),Inches(2.3),CW)
        tb(s,x+Inches(0.2),y+Inches(0.15),Inches(3.5),Inches(0.35),t,15,C1,True)
        tb(s,x+Inches(0.2),y+Inches(0.6),Inches(3.5),Inches(1.5),d,12,CG)

    # 추진 일정
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"6. 추진 일정")
    tbl(s,Inches(0.7),Inches(1.2),[["단계","기간","주요 활동","산출물"],["제안/착수","1주","킥오프, PMP, WBS","헌장, PMP, WBS"],["분석","3주","요구사항, 프로세스","SRS, RTM, BPD"],["설계","3주","아키텍처, ERD, API","설계서, ERD, API 명세"],["구현","5주","서비스 개발, 단위 테스트","소스코드, 테스트 결과"],["테스트","2주","통합/성능/보안/UAT","테스트 결과서"],["이행","1주","배포, 교육, 검수","매뉴얼, 감수 조서"]],[2,1.5,4,4.5])

    # 기대 효과
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"8. 기대 효과")
    effects=[("실험 기록 디지털화","종이 노트 대비 검색/추적 용이, 버전 관리 자동화"),("규정 준수 강화","SHA-256 전자서명으로 위변조 방지, 완전한 감사 추적"),("업무 효율 향상","실시간 협업, 통합 검색, 자동 알림으로 연구 시간 절약"),("자원 관리 최적화","시약 재고 실시간 파악, 장비 예약 충돌 방지")]
    for i,(t,d) in enumerate(effects):
        y=Inches(1.2)+Inches(i*1.4)
        rect(s,Inches(0.5),y,Inches(12.3),Inches(1.1),RGBColor(0xEF,0xF6,0xFF))
        tb(s,Inches(0.7),y+Inches(0.1),Inches(3),Inches(0.3),t,14,C1,True)
        tb(s,Inches(0.7),y+Inches(0.5),Inches(11.5),Inches(0.4),d,12,CD)

    sv(prs,"02_제안서.pptx")

# ════════════════════════════════════════════
# 06. 주간/월간 보고서
# ════════════════════════════════════════════
def g06():
    prs=mkprs(); cover(prs,"주간 업무 보고서","Weekly Progress Report")

    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"1. 프로젝트 현황 요약")
    tbl(s,Inches(0.7),Inches(1.2),[["항목","내용"],["보고 기간",""],["전체 진척률","100%"],["금주 주요 성과","전체 서비스 구현 완료, 통합 테스트 PASS"],["차주 계획","성능 테스트, UAT 수행"],["이슈/리스크","PDF 내보내기 성능 최적화 필요"]],[3,9])

    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"2. 단계별 진척 현황")
    phases=[("제안/착수","100%","완료"),("분석","100%","완료"),("설계","100%","완료"),("구현","100%","완료"),("테스트","100%","완료"),("이행","100%","완료")]
    tbl(s,Inches(0.7),Inches(1.2),[["단계","진척률","상태"]]+[[p,r,s2] for p,r,s2 in phases],[4,2,3])

    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"3. 이슈 및 리스크")
    tbl(s,Inches(0.7),Inches(1.2),[["No","이슈","영향","상태","대응"],["1","PDF 내보내기 28.5초 (목표 30초 근접)","중","진행중","Puppeteer 풀 재사용 적용"],["2","HTTP 보안 헤더 미적용","하","진행중","@fastify/helmet 적용 예정"],["3","비밀번호 복잡도 미적용","중","완료","Zod 정책 추가 완료"]],[1,4,1.5,1.5,4])

    sv(prs,"06_주간월간보고서.pptx")

# ════════════════════════════════════════════
# 21. 설계단계 완료보고서
# ════════════════════════════════════════════
def g21():
    prs=mkprs(); cover(prs,"설계단계 완료보고서","Design Phase Completion Report")

    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"1. 설계 결과 요약")
    tbl(s,Inches(0.7),Inches(1.2),[["항목","내용"],["설계 기간","3주"],["산출물","아키텍처 설계서, ERD, API 명세서(8개 서비스 OpenAPI), 인터페이스 정의서, 프로그램 목록, 프로세스 정의서"],["서비스 수","9개 백엔드 + 2개 프론트엔드"],["DB 스키마","7개 (auth, eln, signature, inventory, scheduler, search, file)"],["API 엔드포인트","144개 (전체 서비스 합산)"],["통신 패턴","HTTP(동기), Redis Stream(비동기), Pub/Sub(실시간), BullMQ(큐), WebSocket"]],[3,9])

    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"2. 아키텍처 설계")
    tb(s,Inches(0.7),Inches(1.1),Inches(11),Inches(0.4),"MSA 기반 9개 서비스, API Gateway 중앙 라우팅, Redis 기반 이벤트/캐시/큐",13,CG)
    svcs=[("gateway","JWT,프록시,SSE"),("auth","인증,RBAC"),("eln","노트,템플릿"),("sig-audit","서명,감사,PDF"),("inventory","시약/장비"),("scheduler","예약"),("search","통합검색"),("file","파일관리"),("collab","실시간협업")]
    for i,(n,d) in enumerate(svcs):
        x=Inches(0.3)+Inches(i*1.43); y=Inches(1.8)
        rect(s,x,y,Inches(1.3),Inches(1.2),RGBColor(0xEF,0xF6,0xFF))
        tb(s,x+Inches(0.05),y+Inches(0.1),Inches(1.2),Inches(0.4),n,10,C1,True,PP_ALIGN.CENTER)
        tb(s,x+Inches(0.05),y+Inches(0.55),Inches(1.2),Inches(0.4),d,8,CG,False,PP_ALIGN.CENTER)

    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"3. 데이터 모델 요약")
    tbl(s,Inches(0.7),Inches(1.2),[["스키마","주요 테이블","레코드 관계"],["auth","Organization, Role, User, Team, TeamMember","Org→User, Role→User, Team↔User"],["eln","Note, NoteRevision, Attachment, NoteLink, Template, Code","Note→Revision, Note→Attachment, Note→Link"],["signature","Signature, AuditLog, ExportJob, Notification","Signature→Note (해시체인)"],["inventory","InventoryItem, InventoryHistory, Category","Item→History"],["scheduler","Resource, Booking","Resource→Booking"],["search","SearchHistory, Favorite, SearchKeywordFavorite","독립"],["file","File, ExportJob","ExportJob→File"]],[2,5,5])

    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"4. 설계 주요 결정")
    items=["Fastify 전 서비스 통일 (Express 대비 2~3배 처리량)","단일 PostgreSQL + Prisma 서비스별 스키마 분리 (운영 단순화)","Redis Stream for 서명→상태전환 (HTTP 폴백 포함, 최소 1회 보장)","BullMQ for PDF 내보내기 (non-blocking, 재시도 3회)","OpenSearch + nori for 한국어 통합 검색","WebSocket + Redis Pub/Sub for 다중 인스턴스 실시간 협업","JWT 듀얼 검증 (jose JWKS + 로컬 시크릿, SSO 선택 지원)"]
    bul(s,Inches(0.7),Inches(1.2),Inches(12),Inches(5),items,14)

    sv(prs,"21_설계단계_완료보고서.pptx")

# ════════════════════════════════════════════
# 40. 교육 계획서/교육자료
# ════════════════════════════════════════════
def g40():
    prs=mkprs(); cover(prs,"LabNote ELN 사용자 교육","User Training Material")

    # 교육 계획
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"교육 계획")
    tbl(s,Inches(0.7),Inches(1.2),[["구분","대상","시간","내용"],["연구원 교육","Researcher 3명","2시간","로그인, 노트 작성/편집, 첨부파일, 인벤토리, 검색, 예약"],["검토자 교육","Reviewer 1명","1시간","노트 검토, 전자서명, PDF 내보내기, 감사 로그"],["관리자 교육","Admin 1명","1.5시간","사용자/팀/역할 관리, 잠금 해제, 대시보드, 설정"],["운영자 교육","IT 운영 1명","1시간","Docker 관리, 헬스체크, DB 백업, 장애 대응"]],[2.5,2.5,1.5,6])

    # 시스템 접속
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"1. 시스템 접속")
    items=["브라우저에서 http://localhost:5173 접속","이메일 / 비밀번호 입력 후 로그인","우측 상단: 언어 전환(한/영), 테마(라이트/다크), 알림 벨, 프로필","좌측 사이드바: 메인 메뉴 6개 + 규정준수 3개 + 관리자 4개"]
    bul(s,Inches(0.7),Inches(1.2),Inches(12),Inches(5),items,15)

    # 연구노트
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"2. 연구노트 작성")
    items=["\"새 노트\" 버튼 → 제목/내용 입력 → 저장 (draft 상태)","편집기 4개 탭: 편집기 | 첨부파일 | 연결(시약/장비) | 버전 이력","상태 전환: draft ↔ in_progress → locked/signed","실시간 협업: 같은 노트를 여러 명이 동시 편집 가능","\"진행 시작\" 버튼: draft → in_progress","수정할 때마다 자동으로 리비전(버전) 생성"]
    bul(s,Inches(0.7),Inches(1.2),Inches(12),Inches(5),items,14)

    # 전자서명
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"3. 전자서명")
    items=["노트 상태가 \"진행 중(in_progress)\"일 때 서명 가능","Reviewer/Admin만 서명 가능 (Researcher는 불가)","편집기에서 \"서명\" 버튼 → 비밀번호 입력 → 서명 완료","서명 완료(signed) = 영구 불변 (수정/삭제/상태변경 모두 불가)","SHA-256 해시체인으로 무결성 보장"]
    bul(s,Inches(0.7),Inches(1.2),Inches(12),Inches(5),items,14)

    # 인벤토리
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"4. 인벤토리 관리")
    items=["유형 탭: 시약/샘플/장비/소모품 등","\"항목 추가\": 이름, 유형, 수량, 바코드, 유효기간 등","수량 조정: 입고(in)/출고(out)/조정(adjust) + 사유","알림: 재고 부족 (수량 ≤ 최소수량), 만료 임박 (N일 전)","행 클릭 → 상세 정보 + 변경 이력 확인"]
    bul(s,Inches(0.7),Inches(1.2),Inches(12),Inches(5),items,14)

    # 예약
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"5. 장비/회의실 예약")
    items=["주간 캘린더 뷰: 월~일, 08:00~17:00","\"새 예약\": 자원 선택 → 제목/날짜/시간 입력","같은 시간대 충돌 시 자동 오류 표시","예약 상태: PENDING → APPROVED/REJECTED → COMPLETED/CANCELLED","승인: Admin 또는 자원 담당자가 승인/반려"]
    bul(s,Inches(0.7),Inches(1.2),Inches(12),Inches(5),items,14)

    # 검색/내보내기
    s=prs.slides.add_slide(prs.slide_layouts[6]); hdr(s,"6. 검색 & 내보내기")
    items=["통합 검색: 노트/템플릿/인벤토리 한 번에 검색","한국어 형태소 분석: \"실험\" → \"실험법\", \"실험 결과\" 매칭","즐겨찾기: 검색어 저장, 문서 북마크","PDF 내보내기: signed/locked 노트 → PDF 자동 생성","ZIP 내보내기: 여러 노트를 한 번에 묶어 다운로드","실시간 진행 상태 표시 (SSE)"]
    bul(s,Inches(0.7),Inches(1.2),Inches(12),Inches(5),items,14)

    # Q&A
    s=prs.slides.add_slide(prs.slide_layouts[6])
    r1=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(13.333),Inches(7.5))
    r1.fill.gradient(); r1.fill.gradient_stops[0].color.rgb=C1; r1.fill.gradient_stops[1].color.rgb=RGBColor(0x0E,0x2A,0x6B); r1.line.fill.background()
    tb(s,Inches(1.5),Inches(2.5),Inches(10),Inches(1),"Q & A",44,CW,True,PP_ALIGN.CENTER)
    tb(s,Inches(1.5),Inches(3.8),Inches(10),Inches(0.5),"문의사항은 시스템 관리자에게 연락해 주세요.",16,RGBColor(0x93,0xC5,0xFD),False,PP_ALIGN.CENTER)

    sv(prs,"40_교육_계획서.pptx")

# ── 실행 ──
print("=== PPTX 생성 시작 ===")
g02(); g06(); g21(); g40()
print("=== PPTX 4건 완료 ===")
