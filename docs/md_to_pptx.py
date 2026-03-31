"""
sequence-diagrams.md → PPT 변환 스크립트
mermaid 다이어그램을 PNG로 렌더링한 뒤 python-pptx로 삽입
"""
import re
import os
import subprocess
import tempfile
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = Path(__file__).parent
MD_FILE = SCRIPT_DIR / "sequence-diagrams.md"
OUTPUT_FILE = SCRIPT_DIR / "sequence-diagrams.pptx"

# 슬라이드 크기: 와이드스크린 16:9
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

TITLE_COLOR = RGBColor(0x2F, 0x54, 0x96)
SUBTITLE_COLOR = RGBColor(0x59, 0x56, 0x59)


def extract_sections(md_text: str):
    """마크다운에서 섹션(제목 + mermaid 블록들)을 추출"""
    sections = []
    lines = md_text.split('\n')

    current_h2 = None
    current_h3 = None
    current_mermaid = None
    in_mermaid = False
    mermaid_lines = []

    for line in lines:
        # H2 헤더
        h2_match = re.match(r'^## (.+)', line)
        if h2_match:
            current_h2 = h2_match.group(1).strip()
            current_h3 = None
            continue

        # H3 헤더
        h3_match = re.match(r'^### (.+)', line)
        if h3_match:
            current_h3 = h3_match.group(1).strip()
            continue

        # mermaid 블록 시작
        if line.strip() == '```mermaid':
            in_mermaid = True
            mermaid_lines = []
            continue

        # mermaid 블록 끝
        if in_mermaid and line.strip() == '```':
            in_mermaid = False
            title = current_h3 or current_h2 or "다이어그램"
            parent_title = current_h2 if current_h3 else None
            sections.append({
                'title': title,
                'parent_title': parent_title,
                'mermaid': '\n'.join(mermaid_lines),
            })
            continue

        if in_mermaid:
            mermaid_lines.append(line)

    # 부록: 이벤트 흐름 매핑 표 추출
    table_match = re.search(
        r'## 부록: 이벤트 흐름 매핑\n\n((?:\|.+\n)+)', md_text
    )
    if table_match:
        sections.append({
            'title': '부록: 이벤트 흐름 매핑',
            'parent_title': None,
            'mermaid': None,
            'table_text': table_match.group(1).strip(),
        })

    return sections


def sanitize_mermaid(code: str) -> str:
    """mermaid 파서가 깨지는 특수문자를 처리"""
    lines = code.split('\n')
    result = []
    for line in lines:
        # 메시지 텍스트 내 세미콜론을 #semi;로 교체 (mermaid HTML entity)
        # 라인에서 ': ' 뒤의 텍스트 부분만 처리
        if ':' in line and ('->>' in line or '-->>' in line or '--)' in line):
            colon_idx = line.index(':', line.index('>>') if '>>' in line else line.index(')'))
            prefix = line[:colon_idx + 1]
            message = line[colon_idx + 1:]
            message = message.replace(';', '#semi;')
            line = prefix + message
        elif line.strip().startswith('Note') and ':' in line:
            colon_idx = line.index(':')
            prefix = line[:colon_idx + 1]
            message = line[colon_idx + 1:]
            message = message.replace(';', '#semi;')
            line = prefix + message
        result.append(line)
    return '\n'.join(result)


def render_mermaid_to_png(mermaid_code: str, output_path: str):
    """mmdc를 사용하여 mermaid → PNG 렌더링"""
    mermaid_code = sanitize_mermaid(mermaid_code)
    with tempfile.NamedTemporaryFile(mode='w', suffix='.mmd', delete=False, encoding='utf-8') as f:
        f.write(mermaid_code)
        mmd_file = f.name

    # mmdc 설정: 높은 해상도
    config = {
        "theme": "default",
        "themeVariables": {
            "fontSize": "14px"
        }
    }
    config_file = mmd_file + '.json'
    with open(config_file, 'w') as f:
        json.dump(config, f)

    mmdc_cmd = os.path.join(os.environ.get('APPDATA', ''), 'npm', 'mmdc.cmd')
    if not os.path.exists(mmdc_cmd):
        mmdc_cmd = 'mmdc'

    try:
        result = subprocess.run(
            [mmdc_cmd, '-i', mmd_file, '-o', output_path,
             '-c', config_file, '-s', '3', '-b', 'white',
             '-w', '2400'],
            capture_output=True, text=True, timeout=60, shell=True
        )
        if result.returncode != 0:
            print(f"  mmdc 경고: {result.stderr[:200]}")
        return os.path.exists(output_path)
    except Exception as e:
        print(f"  mmdc 오류: {e}")
        return False
    finally:
        os.unlink(mmd_file)
        if os.path.exists(config_file):
            os.unlink(config_file)


def add_title_slide(prs: Presentation):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 제목
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "LabNote ELN - 시퀀스 다이어그램"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 부제
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "전자연구노트(ELN) 협업 플랫폼의 핵심 흐름 시퀀스 다이어그램"
    p2.font.size = Pt(18)
    p2.font.color.rgb = SUBTITLE_COLOR
    p2.alignment = PP_ALIGN.CENTER


def add_toc_slide(prs: Presentation, sections):
    """목차 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "목차"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # 고유 H2 제목만 추출
    seen = set()
    h2_titles = []
    for s in sections:
        t = s.get('parent_title') or s['title']
        if t not in seen and not t.startswith('부록'):
            seen.add(t)
            h2_titles.append(t)

    txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, title in enumerate(h2_titles):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = f"{i+1}.  {title}"
        p.font.size = Pt(18)
        p.space_after = Pt(12)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def add_diagram_slide(prs: Presentation, title: str, parent_title: str, png_path: str):
    """다이어그램 이미지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 섹션 제목
    display_title = title
    if parent_title and parent_title != title:
        display_title = f"{parent_title} — {title}"

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = display_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # 이미지 삽입 (가용 영역에 맞추기)
    from PIL import Image
    img = Image.open(png_path)
    img_w, img_h = img.size

    max_w = Inches(12.3)
    max_h = Inches(6.3)
    img_ratio = img_w / img_h
    box_ratio = max_w / max_h

    if img_ratio > box_ratio:
        w = max_w
        h = int(w / img_ratio)
    else:
        h = max_h
        w = int(h * img_ratio)

    left = int((SLIDE_WIDTH - w) / 2)
    top = Inches(1.0)

    slide.shapes.add_picture(png_path, left, top, w, h)


def add_table_slide(prs: Presentation, title: str, table_text: str):
    """마크다운 표를 슬라이드 표로 변환"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    rows_raw = [r.strip() for r in table_text.strip().split('\n') if r.strip()]
    # 헤더, 구분선, 데이터
    header = [c.strip() for c in rows_raw[0].split('|')[1:-1]]
    data_rows = []
    for row in rows_raw[2:]:
        cells = [c.strip() for c in row.split('|')[1:-1]]
        data_rows.append(cells)

    n_rows = len(data_rows) + 1
    n_cols = len(header)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.2),
        Inches(12), Inches(0.4 * n_rows)
    )
    table = table_shape.table

    for i, h in enumerate(header):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True

    for r, row_data in enumerate(data_rows):
        for c, val in enumerate(row_data):
            if c < n_cols:
                cell = table.cell(r + 1, c)
                cell.text = val.replace('`', '')
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)


def main():
    print("=" * 50)
    print("시퀀스 다이어그램 → PPT 변환")
    print("=" * 50)

    md_text = MD_FILE.read_text(encoding='utf-8')
    sections = extract_sections(md_text)
    print(f"\n총 {len(sections)}개 섹션 발견")

    # mermaid → PNG 렌더링
    tmp_dir = tempfile.mkdtemp(prefix='mermaid_')
    print(f"임시 디렉토리: {tmp_dir}\n")

    for i, sec in enumerate(sections):
        if sec.get('mermaid'):
            png_path = os.path.join(tmp_dir, f"diagram_{i}.png")
            print(f"[{i+1}/{len(sections)}] 렌더링: {sec['title'][:40]}...")
            ok = render_mermaid_to_png(sec['mermaid'], png_path)
            sec['png_path'] = png_path if ok else None
            print(f"  → {'성공' if ok else '실패'}")
        else:
            sec['png_path'] = None

    # PPT 생성
    print("\nPPT 생성 중...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_title_slide(prs)
    add_toc_slide(prs, sections)

    for sec in sections:
        if sec.get('png_path') and os.path.exists(sec['png_path']):
            add_diagram_slide(
                prs, sec['title'],
                sec.get('parent_title'),
                sec['png_path']
            )
        elif sec.get('table_text'):
            add_table_slide(prs, sec['title'], sec['table_text'])
        else:
            print(f"  건너뜀 (이미지 없음): {sec['title']}")

    prs.save(str(OUTPUT_FILE))
    print(f"\n완료: {OUTPUT_FILE}")
    print(f"슬라이드 수: {len(prs.slides)}")

    # 임시 파일 정리
    import shutil
    shutil.rmtree(tmp_dir, ignore_errors=True)


if __name__ == '__main__':
    main()
