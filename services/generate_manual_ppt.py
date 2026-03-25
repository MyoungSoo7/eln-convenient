#!/usr/bin/env python3
"""LabNote ELN 사용자 매뉴얼 PPT 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
C_PRIMARY    = RGBColor(0x1E, 0x40, 0xAF)  # 진한 파랑
C_SECONDARY  = RGBColor(0x3B, 0x82, 0xF6)  # 밝은 파랑
C_ACCENT     = RGBColor(0x06, 0xB6, 0xD4)  # 시안
C_SUCCESS    = RGBColor(0x10, 0xB9, 0x81)  # 녹색
C_WARNING    = RGBColor(0xF5, 0x9E, 0x0B)  # 주황
C_DANGER     = RGBColor(0xEF, 0x44, 0x44)  # 빨강
C_DARK       = RGBColor(0x1E, 0x29, 0x3B)  # 거의 검정
C_GRAY       = RGBColor(0x64, 0x74, 0x8B)  # 회색
C_LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)  # 밝은 배경
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_BLUE    = RGBColor(0xEF, 0xF6, 0xFF)  # 연한 파랑 배경
C_BG_CYAN    = RGBColor(0xEC, 0xFD, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── 유틸리티 함수 ──

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gradient_rect(slide, left, top, width, height, color1, color2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2
    return shape

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text="", font_size=14, color=C_DARK,
                 bold=False, align=PP_ALIGN.LEFT, font_name="맑은 고딕", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    txBox.text_frame.paragraphs[0].space_after = Pt(0)
    txBox.text_frame.paragraphs[0].space_before = Pt(0)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=C_DARK,
                    spacing=Pt(6), bullet_char="\u2022", font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.space_before = Pt(0)
    return txBox

def add_card(slide, left, top, width, height, title, body_lines, title_color=C_PRIMARY,
             bg_color=C_WHITE, border_color=None):
    card = add_rect(slide, left, top, width, height, fill_color=bg_color, line_color=border_color or RGBColor(0xE2,0xE8,0xF0))
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.35),
                 title, font_size=13, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_lines, font_size=11, color=C_DARK, spacing=Pt(4))
    return card

def add_table_slide(slide, left, top, rows_data, col_widths, font_size=11, header_color=C_PRIMARY):
    """rows_data: list of lists. First row = header."""
    rows = len(rows_data)
    cols = len(rows_data[0])
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(rows, cols, left, top, Inches(total_w), Inches(rows * 0.38))
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "맑은 고딕"
                paragraph.space_after = Pt(0)
                paragraph.space_before = Pt(0)
                if ri == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = C_WHITE
                else:
                    paragraph.font.color.rgb = C_DARK
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE
    return table_shape

def section_header_bar(slide, text):
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.0), fill_color=C_PRIMARY)
    add_text_box(slide, Inches(0.7), Inches(0.2), Inches(10), Inches(0.6),
                 text, font_size=28, color=C_WHITE, bold=True)

def slide_title_bar(slide, section_num, title, subtitle=""):
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.85), fill_color=C_PRIMARY)
    num_txt = f"{section_num}.  " if section_num else ""
    add_text_box(slide, Inches(0.7), Inches(0.12), Inches(11), Inches(0.55),
                 f"{num_txt}{title}", font_size=24, color=C_WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.7), Inches(0.55), Inches(11), Inches(0.3),
                     subtitle, font_size=13, color=RGBColor(0xBF,0xDB,0xFE))

def page_number(slide, num):
    add_text_box(slide, Inches(12.4), Inches(7.1), Inches(0.8), Inches(0.3),
                 str(num), font_size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)

pn = [0]
def next_page():
    pn[0] += 1
    return pn[0]

# ════════════════════════════════════════════════════════════════
# SLIDE 1: 표지
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_gradient_rect(slide, 0, 0, W, H, C_PRIMARY, RGBColor(0x0E, 0x2A, 0x6B))

# 제목
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
             "LabNote ELN", font_size=52, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.85), Inches(10), Inches(0.6),
             "전자연구노트 협업 플랫폼", font_size=28, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

# 구분선
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.7), Inches(2.3), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = C_ACCENT
line.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(0.6),
             "사 용 자  매 뉴 얼  v1.0", font_size=24, color=C_WHITE, align=PP_ALIGN.CENTER)

# 하단 키워드
keywords = ["연구노트  |  전자서명  |  프로토콜  |  인벤토리  |  예약  |  통합검색  |  감사로그  |  내보내기"]
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.4),
             keywords[0], font_size=13, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 2: 목차
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "", "목차", "Table of Contents")
p = next_page()

toc_left = [
    ("01", "시스템 개요"),
    ("02", "접속 및 로그인"),
    ("03", "화면 구성"),
    ("04", "대시보드"),
    ("05", "연구노트"),
    ("06", "프로토콜 (템플릿)"),
    ("07", "인벤토리"),
]
toc_right = [
    ("08", "장비/회의실 예약"),
    ("09", "통합 검색"),
    ("10", "전자서명 / 규정 준수"),
    ("11", "감사 로그"),
    ("12", "내보내기 (PDF/ZIP)"),
    ("13", "알림"),
    ("14", "관리자 기능 / 역할 / 권한"),
]

for i, (num, title) in enumerate(toc_left):
    y = Inches(1.3) + Inches(i * 0.65)
    add_rect(slide, Inches(1.0), y, Inches(0.6), Inches(0.45), fill_color=C_PRIMARY)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(0.6), Inches(0.35),
                 num, font_size=16, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.8), y + Inches(0.05), Inches(4), Inches(0.35),
                 title, font_size=16, color=C_DARK, bold=False)

for i, (num, title) in enumerate(toc_right):
    y = Inches(1.3) + Inches(i * 0.65)
    add_rect(slide, Inches(7.0), y, Inches(0.6), Inches(0.45), fill_color=C_PRIMARY)
    add_text_box(slide, Inches(7.0), y + Inches(0.05), Inches(0.6), Inches(0.35),
                 num, font_size=16, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.8), y + Inches(0.05), Inches(4.5), Inches(0.35),
                 title, font_size=16, color=C_DARK, bold=False)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 3: 시스템 개요
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "01", "시스템 개요", "System Overview")
p = next_page()

add_text_box(slide, Inches(0.7), Inches(1.1), Inches(11), Inches(0.4),
             "LabNote ELN은 연구실에서 사용하는 전자연구노트 협업 플랫폼입니다.",
             font_size=15, color=C_DARK)

features = [
    ("연구노트", "노트 작성, 편집, 버전 관리\n실시간 다중 사용자 협업", C_PRIMARY),
    ("전자서명", "SHA-256 해시체인 서명\n무결성 검증, 규정 준수", C_SUCCESS),
    ("프로토콜", "실험법/SOP 템플릿 관리\n공유 및 재사용", C_SECONDARY),
    ("인벤토리", "시약/샘플/장비 재고 추적\n만료/부족 알림", C_WARNING),
    ("예약", "장비/회의실 예약\n승인/반려 워크플로", C_ACCENT),
    ("통합 검색", "노트/템플릿/인벤토리 통합\n한국어 형태소 분석", RGBColor(0x8B,0x5C,0xF6)),
    ("감사 로그", "모든 활동 자동 기록\n필터링 및 추적", C_DANGER),
    ("내보내기", "PDF / ZIP / 보고서\n비동기 생성, 실시간 알림", RGBColor(0xEC,0x48,0x99)),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.7) + Inches(col * 3.1)
    y = Inches(1.8) + Inches(row * 2.5)

    card = add_rect(slide, x, y, Inches(2.8), Inches(2.1), fill_color=C_WHITE,
                    line_color=RGBColor(0xE2,0xE8,0xF0))

    # color bar at top
    bar = add_rect(slide, x, y, Inches(2.8), Inches(0.08), fill_color=color)

    add_text_box(slide, x + Inches(0.2), y + Inches(0.25), Inches(2.4), Inches(0.35),
                 title, font_size=16, color=color, bold=True)
    lines = desc.split('\n')
    for li, line_text in enumerate(lines):
        add_text_box(slide, x + Inches(0.2), y + Inches(0.7) + Inches(li * 0.35), Inches(2.4), Inches(0.3),
                     line_text, font_size=12, color=C_GRAY)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 4: 시스템 아키텍처
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "01", "시스템 아키텍처", "Architecture")
p = next_page()

add_text_box(slide, Inches(0.7), Inches(1.1), Inches(11), Inches(0.4),
             "MSA 기반 온프레미스 Docker Compose 배포 구조",
             font_size=14, color=C_GRAY)

# Gateway box
add_rect(slide, Inches(1.5), Inches(1.7), Inches(10.3), Inches(0.7), fill_color=C_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.5),
             "API Gateway  (Fastify :8000)  —  JWT 검증  |  프록시  |  Rate Limit  |  SSE  |  대시보드 집계",
             font_size=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

services_row = [
    ("auth\n:8001", "조직/사용자\nRBAC, JWT"),
    ("eln\n:8002", "연구노트\n버전관리"),
    ("sig-audit\n:8003", "전자서명\n감사로그"),
    ("inventory\n:8004", "시약/장비\n재고관리"),
    ("scheduler\n:8005", "예약\n승인 워크플로"),
    ("search\n:8006", "통합검색\nOpenSearch"),
    ("file\n:8008", "파일 업/다운\nMinIO"),
    ("collab\n:8009", "실시간 협업\nWebSocket"),
]

for i, (name, desc) in enumerate(services_row):
    x = Inches(0.5) + Inches(i * 1.55)
    y = Inches(2.8)
    add_rect(slide, x, y, Inches(1.4), Inches(1.5), fill_color=C_BG_BLUE, line_color=C_SECONDARY)
    add_text_box(slide, x + Inches(0.05), y + Inches(0.1), Inches(1.3), Inches(0.6),
                 name, font_size=11, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.05), y + Inches(0.7), Inches(1.3), Inches(0.7),
                 desc, font_size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

# Infra row
infra = [
    ("PostgreSQL 15", "단일 인스턴스\n서비스별 스키마 분리"),
    ("Redis 7", "캐시, Stream, Pub/Sub\nBullMQ, 세션 관리"),
    ("MinIO", "S3 호환 오브젝트 스토리지\n파일 업로드/다운로드"),
    ("OpenSearch 2", "통합 검색 엔진\n한국어 nori 분석기"),
    ("Keycloak 24", "SSO (선택)\nJWKS 연동"),
]

add_text_box(slide, Inches(0.7), Inches(4.6), Inches(4), Inches(0.3),
             "인프라 컴포넌트", font_size=14, color=C_DARK, bold=True)

for i, (name, desc) in enumerate(infra):
    x = Inches(0.5) + Inches(i * 2.5)
    y = Inches(5.05)
    add_rect(slide, x, y, Inches(2.3), Inches(1.3), fill_color=RGBColor(0xF0,0xFD,0xF4), line_color=C_SUCCESS)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.1), Inches(2.1), Inches(0.35),
                 name, font_size=12, color=C_SUCCESS, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.5), Inches(2.1), Inches(0.7),
                 desc, font_size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 5: 접속 및 로그인
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "02", "접속 및 로그인", "Login & Settings")
p = next_page()

# 로그인 섹션
add_text_box(slide, Inches(0.7), Inches(1.1), Inches(5), Inches(0.4),
             "로그인 절차", font_size=18, color=C_PRIMARY, bold=True)

login_steps = [
    "브라우저에서 http://localhost:5173 접속",
    "이메일과 비밀번호 입력",
    "\"로그인\" 버튼 클릭",
    "로그인 성공 시 대시보드(/) 자동 이동",
]
add_bullet_list(slide, Inches(0.9), Inches(1.6), Inches(5), Inches(2.0),
                login_steps, font_size=13, spacing=Pt(8))

# 로그아웃
add_text_box(slide, Inches(0.7), Inches(3.5), Inches(5), Inches(0.4),
             "로그아웃", font_size=18, color=C_PRIMARY, bold=True)
add_bullet_list(slide, Inches(0.9), Inches(4.0), Inches(5), Inches(0.8),
                ["우측 상단 사용자 아바타 클릭 → 로그아웃 선택"],
                font_size=13, spacing=Pt(8))

# 우측: 설정
add_text_box(slide, Inches(7.0), Inches(1.1), Inches(5), Inches(0.4),
             "환경 설정", font_size=18, color=C_PRIMARY, bold=True)

settings = [
    "언어 전환: 헤더의 언어(🌐) 버튼 → 한국어/영어 전환",
    "테마 전환: 헤더의 테마(🌙) 버튼 → 라이트/다크 모드",
    "로그인 화면에서도 언어 전환 가능",
    "언어 설정은 브라우저에 저장되어 유지",
]
add_bullet_list(slide, Inches(7.2), Inches(1.6), Inches(5.5), Inches(2.0),
                settings, font_size=13, spacing=Pt(8))

# 접속 URL 테이블
add_text_box(slide, Inches(7.0), Inches(3.5), Inches(5), Inches(0.4),
             "접속 URL", font_size=18, color=C_PRIMARY, bold=True)
add_table_slide(slide, Inches(7.0), Inches(4.0), [
    ["용도", "URL"],
    ["메인 프론트엔드", "http://localhost:5173"],
    ["인벤토리 전용", "http://localhost:3000"],
    ["API Gateway", "http://localhost:8000"],
    ["MinIO 콘솔", "http://localhost:9001"],
], [2.5, 3.5], font_size=11)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 6: 화면 구성
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "03", "화면 구성", "UI Layout & Navigation")
p = next_page()

# 레이아웃 다이어그램
add_rect(slide, Inches(0.7), Inches(1.2), Inches(5.5), Inches(0.6), fill_color=C_PRIMARY)
add_text_box(slide, Inches(0.7), Inches(1.25), Inches(5.5), Inches(0.5),
             "  [=]  LabNote                                🌐  🌙  🔔  👤  헤더 영역",
             font_size=12, color=C_WHITE, bold=False)

add_rect(slide, Inches(0.7), Inches(1.85), Inches(1.2), Inches(4.0), fill_color=C_BG_BLUE, line_color=C_SECONDARY)
add_text_box(slide, Inches(0.75), Inches(1.95), Inches(1.1), Inches(0.3),
             "사이드바", font_size=11, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
sb_items = ["대시보드", "연구노트", "프로토콜", "인벤토리", "예약", "검색", "─────", "서명 현황", "감사 로그", "내보내기", "─────", "관리자 ▼"]
for i, item in enumerate(sb_items):
    add_text_box(slide, Inches(0.8), Inches(2.3) + Inches(i * 0.28), Inches(1.0), Inches(0.25),
                 item, font_size=9, color=C_GRAY if "──" in item else C_DARK)

add_rect(slide, Inches(1.95), Inches(1.85), Inches(4.25), Inches(4.0), fill_color=C_LIGHT_GRAY, line_color=RGBColor(0xE2,0xE8,0xF0))
add_text_box(slide, Inches(1.95), Inches(3.4), Inches(4.25), Inches(0.5),
             "메인 콘텐츠 영역", font_size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

# 메뉴 테이블
add_text_box(slide, Inches(7.0), Inches(1.1), Inches(5), Inches(0.4),
             "사이드바 메뉴 구성", font_size=16, color=C_PRIMARY, bold=True)

add_table_slide(slide, Inches(7.0), Inches(1.55), [
    ["그룹", "메뉴", "경로"],
    ["메인", "대시보드", "/"],
    ["메인", "연구노트", "/notes"],
    ["메인", "프로토콜", "/protocols"],
    ["메인", "인벤토리", "/inventory"],
    ["메인", "예약", "/scheduler"],
    ["메인", "검색", "/search"],
    ["규정 준수", "서명 현황", "/signatures"],
    ["규정 준수", "감사 로그", "/audit-logs"],
    ["규정 준수", "내보내기", "/exports"],
    ["관리자", "사용자 관리", "/admin/users"],
    ["관리자", "역할 관리", "/admin/roles"],
    ["관리자", "코드 관리", "/admin/codes"],
    ["관리자", "설정", "/admin/settings"],
], [1.2, 1.6, 2.2], font_size=10)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 7: 대시보드
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "04", "대시보드", "Dashboard")
p = next_page()

add_text_box(slide, Inches(0.7), Inches(1.1), Inches(11), Inches(0.35),
             "로그인 후 첫 화면. 현재 작업 상황을 한눈에 파악할 수 있습니다. (60초마다 자동 갱신)",
             font_size=13, color=C_GRAY)

# 내 대시보드
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5), Inches(0.35),
             "내 대시보드 (모든 사용자)", font_size=16, color=C_PRIMARY, bold=True)

# stat cards
stat_cards = [
    ("전체 노트", "12", C_PRIMARY),
    ("진행 중", "4", C_WARNING),
    ("인벤토리", "28", C_SUCCESS),
    ("미읽은 알림", "3", C_DANGER),
]
for i, (label, val, color) in enumerate(stat_cards):
    x = Inches(0.7) + Inches(i * 1.5)
    add_rect(slide, x, Inches(2.05), Inches(1.35), Inches(0.85), fill_color=C_WHITE, line_color=color)
    add_text_box(slide, x + Inches(0.1), Inches(2.1), Inches(1.15), Inches(0.3),
                 label, font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(2.4), Inches(1.15), Inches(0.4),
                 val, font_size=22, color=color, bold=True, align=PP_ALIGN.CENTER)

personal_items = [
    "최근 노트: 최근 수정한 노트 4건 (클릭 → 노트 이동)",
    "예정 예약: 다가오는 예약 4건 (승인/대기 뱃지)",
    "최근 활동: 최근 감사 로그 4건",
    "알림: 재고 부족/만료 임박 아이템 최대 5건",
]
add_bullet_list(slide, Inches(0.9), Inches(3.1), Inches(5.5), Inches(2.0),
                personal_items, font_size=12, spacing=Pt(6))

# 조직 대시보드
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.35),
             "조직 대시보드 (Admin 전용)", font_size=16, color=C_PRIMARY, bold=True)

org_items = [
    "노트 통계: 상태별(초안/진행/서명/잠김) 분류",
    "규정 준수율: 서명 완료 비율 (complianceRate)",
    "인벤토리: 총 수량, 재고 부족, 만료 임박",
    "예약 현황: 대기 중인 예약 수",
    "최근 조직 활동 로그",
    "재고 부족/만료 임박 아이템 상세",
]
add_bullet_list(slide, Inches(7.2), Inches(2.1), Inches(5.3), Inches(3.5),
                org_items, font_size=12, spacing=Pt(6))

add_rect(slide, Inches(7.0), Inches(4.8), Inches(5.5), Inches(0.7), fill_color=C_BG_BLUE)
add_text_box(slide, Inches(7.2), Inches(4.9), Inches(5.1), Inches(0.5),
             "💡 탭으로 \"내 대시보드\"와 \"조직 대시보드\"를 전환할 수 있습니다.",
             font_size=12, color=C_PRIMARY)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 8: 연구노트 - 목록
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "05", "연구노트 — 목록 및 관리", "Notes List")
p = next_page()

# 좌측: 기능 설명
add_text_box(slide, Inches(0.7), Inches(1.1), Inches(5.5), Inches(0.35),
             "노트 목록 화면 (/notes)", font_size=16, color=C_PRIMARY, bold=True)

note_list_items = [
    "상태 필터: 전체 / 초안 / 진행 중 / 서명 완료 / 잠김",
    "검색: 제목 또는 태그로 필터링",
    "템플릿 필터: 프로토콜 \"노트 보기\" 시 자동 적용",
    "\"새 노트\" 버튼 → 노트 편집기로 이동",
    "노트 카드 클릭 → 해당 노트 편집기로 이동",
]
add_bullet_list(slide, Inches(0.9), Inches(1.6), Inches(5.5), Inches(2.0),
                note_list_items, font_size=13, spacing=Pt(8))

# 상태 전환 테이블
add_text_box(slide, Inches(0.7), Inches(3.5), Inches(5.5), Inches(0.35),
             "목록에서 상태 전환", font_size=14, color=C_DARK, bold=True)
add_table_slide(slide, Inches(0.7), Inches(3.95), [
    ["현재 상태", "변경 가능", "필요 역할"],
    ["초안(draft)", "→ 진행 중", "모든 역할"],
    ["진행 중", "→ 초안", "모든 역할"],
    ["진행 중", "→ 잠김", "Reviewer, Admin"],
    ["잠김(locked)", "→ 초안 (해제)", "Admin (비밀번호 필요)"],
], [1.7, 1.7, 2.2], font_size=11)

# 우측: 삭제/제약
add_text_box(slide, Inches(7.0), Inches(1.1), Inches(5.5), Inches(0.35),
             "노트 생성", font_size=16, color=C_PRIMARY, bold=True)
create_items = [
    "\"새 노트\" 버튼 클릭",
    "제목과 내용을 입력",
    "저장 버튼 클릭 → 초안(draft) 상태로 생성",
    "프로토콜에서 \"노트 작성\" 시 템플릿 내용 자동 채움",
]
add_bullet_list(slide, Inches(7.2), Inches(1.6), Inches(5.3), Inches(1.8),
                create_items, font_size=13, spacing=Pt(8))

add_text_box(slide, Inches(7.0), Inches(3.5), Inches(5.5), Inches(0.35),
             "삭제 규칙", font_size=16, color=C_DANGER, bold=True)
delete_items = [
    "초안(draft) / 진행 중(in_progress)만 삭제 가능",
    "잠김(locked) / 서명 완료(signed)는 삭제 불가",
    "마우스 오버 시 휴지통 아이콘 표시",
    "삭제 전 확인 대화상자 표시",
]
add_bullet_list(slide, Inches(7.2), Inches(4.0), Inches(5.3), Inches(2.0),
                delete_items, font_size=13, spacing=Pt(8))

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: 연구노트 - 편집기
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "05", "연구노트 — 편집기", "Note Editor")
p = next_page()

add_text_box(slide, Inches(0.7), Inches(1.1), Inches(11), Inches(0.35),
             "노트 편집기 (/notes/:id) — 작성, 편집, 실시간 협업, 첨부파일, 연결, 버전 관리를 모두 수행하는 핵심 화면",
             font_size=13, color=C_GRAY)

# 4 tab cards
tabs = [
    ("편집기", [
        "제목 입력 필드 + 본문 텍스트 영역",
        "저장 버튼 (수정 시 리비전 자동 생성)",
        "\"진행 시작\" 버튼 (draft → in_progress)",
        "\"서명\" 버튼 (Reviewer/Admin, 비밀번호 필요)",
        "잠김/서명 완료 시 편집 비활성화",
    ], C_PRIMARY),
    ("첨부파일", [
        "드래그 앤 드롭 / 클릭 파일 업로드",
        "파일 목록: 이름, MIME, 크기, 날짜",
        "클릭으로 다운로드/열기",
        "마우스 오버 시 삭제 아이콘",
        "새 노트(저장 전) / 잠긴 노트에서 비활성화",
    ], C_SUCCESS),
    ("연결 (시약/장비)", [
        "\"시약/장비 연결\" → 인벤토리 검색 대화상자",
        "검색 후 클릭하여 연결",
        "유형 뱃지 + 이름 표시",
        "X 버튼으로 연결 해제",
        "새 노트 / 잠긴 노트에서 비활성화",
    ], C_WARNING),
    ("버전 이력", [
        "탭 최초 클릭 시 이력 로드",
        "타임라인 뷰",
        "버전 번호, 시간, 변경 요약, 변경자",
        "전체 리비전 내용 확인 가능",
        "",
    ], RGBColor(0x8B,0x5C,0xF6)),
]

for i, (tab_name, items, color) in enumerate(tabs):
    x = Inches(0.5) + Inches(i * 3.1)
    y = Inches(1.65)
    add_rect(slide, x, y, Inches(2.9), Inches(0.4), fill_color=color)
    add_text_box(slide, x, y + Inches(0.05), Inches(2.9), Inches(0.3),
                 tab_name, font_size=13, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_bullet_list(slide, x + Inches(0.15), y + Inches(0.5), Inches(2.6), Inches(2.5),
                    [it for it in items if it], font_size=11, color=C_DARK, spacing=Pt(5))

# 실시간 협업 박스
add_rect(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.5), fill_color=C_BG_CYAN, line_color=C_ACCENT)
add_text_box(slide, Inches(0.7), Inches(4.7), Inches(5), Inches(0.3),
             "실시간 협업 (WebSocket)", font_size=14, color=C_ACCENT, bold=True)
collab_items = [
    "같은 노트를 여러 사용자가 동시에 열면 자동 활성화",
    "접속 중인 사용자의 색상 아바타가 상단에 표시",
    "다른 사용자의 편집 내용이 실시간 반영 (400ms 디바운스)",
    "커서 위치(awareness)가 실시간 공유",
]
add_bullet_list(slide, Inches(0.9), Inches(5.05), Inches(11), Inches(1.0),
                collab_items, font_size=12, spacing=Pt(5))

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 10: 프로토콜 (템플릿)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "06", "프로토콜 (템플릿)", "Protocols & Templates")
p = next_page()

add_text_box(slide, Inches(0.7), Inches(1.1), Inches(11), Inches(0.35),
             "반복 사용하는 실험 절차나 SOP를 템플릿으로 관리하고, 이를 기반으로 노트를 빠르게 생성합니다.",
             font_size=13, color=C_GRAY)

# 좌측: 생성
add_text_box(slide, Inches(0.7), Inches(1.65), Inches(5.5), Inches(0.35),
             "프로토콜 생성", font_size=16, color=C_PRIMARY, bold=True)
proto_create = [
    "\"새 템플릿\" 버튼 클릭",
    "제목 (필수) 입력",
    "카테고리 선택 (일반/실험법/분석법/품질관리/SOP/안전절차)",
    "설명, 내용(본문), 태그 입력",
    "공개 여부 토글 (공개 시 다른 조직에서도 조회 가능)",
    "저장 클릭",
]
add_bullet_list(slide, Inches(0.9), Inches(2.1), Inches(5.3), Inches(2.5),
                proto_create, font_size=13, spacing=Pt(7))

# 우측: 활용
add_text_box(slide, Inches(7.0), Inches(1.65), Inches(5.5), Inches(0.35),
             "프로토콜 활용", font_size=16, color=C_PRIMARY, bold=True)

add_table_slide(slide, Inches(7.0), Inches(2.1), [
    ["기능", "설명"],
    ["노트 작성", "프로토콜 카드 → \"노트 작성\" → 템플릿 내용이 미리 채워진 새 노트 생성"],
    ["노트 보기", "\"노트 보기\" → 이 프로토콜로 생성된 노트 목록으로 이동"],
    ["복사", "프로토콜을 복사하여 새 프로토콜 생성 (원본 복사 횟수 증가)"],
    ["편집", "연필 아이콘 → 수정 대화상자 (소유자 / 관리자만)"],
    ["삭제", "확인 후 삭제 (소유자 / 관리자만)"],
], [1.5, 4.5], font_size=11)

# 하단 참고
add_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.0), fill_color=C_BG_BLUE)
add_text_box(slide, Inches(0.7), Inches(5.3), Inches(11.9), Inches(0.3),
             "프로토콜 카드 표시 정보", font_size=13, color=C_PRIMARY, bold=True)
add_text_box(slide, Inches(0.7), Inches(5.65), Inches(11.9), Inches(0.4),
             "카테고리 뱃지  |  제목  |  태그  |  사용 횟수(\"N회 사용됨\")  |  공개/비공개 뱃지  |  설명 미리보기",
             font_size=12, color=C_DARK)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 11: 인벤토리
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "07", "인벤토리 — 시약/장비 관리", "Inventory Management")
p = next_page()

# 좌: 목록 기능
add_text_box(slide, Inches(0.7), Inches(1.1), Inches(5.5), Inches(0.35),
             "목록 및 필터", font_size=16, color=C_PRIMARY, bold=True)

inv_features = [
    "알림 배너: 만료 임박 / 재고 부족 아이템 수 표시 (클릭 시 필터)",
    "유형 탭: 시약, 샘플, 장비, 소모품, 항체, 플라스미드, 세포주 등",
    "상태 필터: 사용 가능/사용 중/소진/만료/폐기/유지보수",
    "텍스트 검색: 이름으로 검색 (300ms 디바운스)",
]
add_bullet_list(slide, Inches(0.9), Inches(1.6), Inches(5.3), Inches(1.5),
                inv_features, font_size=12, spacing=Pt(6))

add_text_box(slide, Inches(0.7), Inches(3.3), Inches(5.5), Inches(0.35),
             "아이템 추가", font_size=16, color=C_PRIMARY, bold=True)
add_table_slide(slide, Inches(0.7), Inches(3.7), [
    ["필드", "필수", "설명"],
    ["이름", "O", "아이템명"],
    ["유형", "O", "시약, 샘플, 장비 등"],
    ["수량 / 단위", "", "초기 수량, mL/g/EA 등"],
    ["최소 수량", "", "이하 시 재고 부족 알림"],
    ["바코드", "", "고유 바코드 번호"],
    ["유효기간", "", "만료일"],
    ["만료 경고 일수", "", "기본 30일"],
    ["태그", "", "분류 태그 (여러 개)"],
], [1.5, 0.5, 3.5], font_size=10)

# 우: 수량 조정 & 상세
add_text_box(slide, Inches(7.0), Inches(1.1), Inches(5.5), Inches(0.35),
             "수량 조정", font_size=16, color=C_PRIMARY, bold=True)
qty_steps = [
    "아이템 행의 수량 조정(↕) 버튼 클릭",
    "변경 유형 선택:",
    "  • 입고(in): 수량 증가",
    "  • 출고(out): 수량 감소",
    "  • 조정(adjust): 절대값 설정",
    "수량 입력 + 사유 입력 (선택)",
    "\"이전 수량 → 변경 후 수량\" 미리보기 확인 후 확인",
]
add_bullet_list(slide, Inches(7.2), Inches(1.6), Inches(5.3), Inches(2.5),
                qty_steps, font_size=12, spacing=Pt(5))

add_text_box(slide, Inches(7.0), Inches(4.0), Inches(5.5), Inches(0.35),
             "아이템 상세 / 이력", font_size=16, color=C_PRIMARY, bold=True)
detail_items = [
    "아이템 행 클릭 → 상세 정보 대화상자",
    "모든 필드 정보 + 변경 이력 테이블:",
    "  날짜, 변경 유형, 이전/이후 수량,",
    "  이전/이후 상태, 사유, 수행자",
]
add_bullet_list(slide, Inches(7.2), Inches(4.5), Inches(5.3), Inches(1.5),
                detail_items, font_size=12, spacing=Pt(5))

# 알림 표시 규칙
add_rect(slide, Inches(7.0), Inches(5.8), Inches(5.5), Inches(0.9), fill_color=RGBColor(0xFF,0xF7,0xED))
add_text_box(slide, Inches(7.2), Inches(5.85), Inches(5.1), Inches(0.8),
             "만료: 빨간색 뱃지 | N일 남음: 주황색 뱃지\n재고 부족 (수량 ≤ 최소수량): 주황색 수량 + 경고 아이콘",
             font_size=11, color=C_WARNING)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 12: 예약
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "08", "장비/회의실 예약", "Scheduler")
p = next_page()

# 좌: 예약 기능
add_text_box(slide, Inches(0.7), Inches(1.1), Inches(5.5), Inches(0.35),
             "예약 생성 및 관리", font_size=16, color=C_PRIMARY, bold=True)

booking_items = [
    "캘린더 뷰: 주간 월~일, 08:00~17:00 시간 슬롯",
    "  승인(APPROVED): 파란색  |  대기(PENDING): 노란색",
    "주 이동: ◀ 이전/다음 ▶, \"오늘\" 버튼",
    "",
    "예약 생성 절차:",
    "  1. \"새 예약\" 버튼 클릭",
    "  2. 자원 선택 (장비 또는 회의실)",
    "  3. 제목/용도, 날짜, 시작/종료 시간 입력",
    "  4. 같은 시간에 충돌 시 오류 표시",
]
add_bullet_list(slide, Inches(0.9), Inches(1.6), Inches(5.3), Inches(3.0),
                [i for i in booking_items if i], font_size=12, spacing=Pt(5))

# 우: 상태 흐름
add_text_box(slide, Inches(7.0), Inches(1.1), Inches(5.5), Inches(0.35),
             "예약 상태 흐름", font_size=16, color=C_PRIMARY, bold=True)

add_table_slide(slide, Inches(7.0), Inches(1.55), [
    ["현재", "→ 전환", "수행자"],
    ["PENDING", "APPROVED", "Admin, 자원 담당자"],
    ["PENDING", "REJECTED", "Admin, 자원 담당자"],
    ["PENDING", "CANCELLED", "예약자 본인, Admin"],
    ["APPROVED", "COMPLETED", "예약자 본인, Admin"],
    ["APPROVED", "CANCELLED", "예약자 본인, Admin"],
], [1.5, 1.8, 2.5], font_size=11)

# 자원 관리
add_text_box(slide, Inches(7.0), Inches(4.0), Inches(5.5), Inches(0.35),
             "자원 관리 (Admin 전용)", font_size=16, color=C_PRIMARY, bold=True)
res_items = [
    "장비(EQUIPMENT) / 회의실(ROOM) 유형 구분",
    "자원 등록: 이름, 유형, 위치, 설명, 수용 인원(회의실)",
    "승인 담당자 지정 가능 (미지정 시 Admin이 승인)",
    "비활성화: 대기/승인 예약이 있으면 비활성화 불가",
]
add_bullet_list(slide, Inches(7.2), Inches(4.5), Inches(5.3), Inches(2.0),
                res_items, font_size=12, spacing=Pt(6))

# 승인 참고 박스
add_rect(slide, Inches(0.5), Inches(5.5), Inches(5.8), Inches(1.2), fill_color=C_BG_BLUE)
add_text_box(slide, Inches(0.7), Inches(5.6), Inches(5.4), Inches(0.3),
             "승인/반려 (Reviewer, Admin, 자원 담당자)", font_size=13, color=C_PRIMARY, bold=True)
add_text_box(slide, Inches(0.7), Inches(5.95), Inches(5.4), Inches(0.7),
             "승인: 클릭 시 즉시 승인 (DB 비관적 락으로 동시 승인 방지)\n반려: 사유 입력 대화상자 → 반려 처리",
             font_size=12, color=C_DARK)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 13: 통합 검색
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "09", "통합 검색", "Unified Search")
p = next_page()

add_text_box(slide, Inches(0.7), Inches(1.1), Inches(11), Inches(0.35),
             "노트, 템플릿, 인벤토리를 하나의 검색창에서 통합 검색합니다. 한국어 형태소 분석(nori) 지원.",
             font_size=13, color=C_GRAY)

# 검색
add_text_box(slide, Inches(0.7), Inches(1.65), Inches(5.5), Inches(0.35),
             "검색 기능", font_size=16, color=C_PRIMARY, bold=True)
search_items = [
    "검색창에 키워드 입력 (300ms 디바운스 자동 검색)",
    "결과 탭으로 구분: 노트 / 템플릿 / 인벤토리",
    "각 탭에 결과 건수 표시",
    "결과 카드 클릭 → 해당 항목으로 이동",
    "한국어 형태소 분석: \"실험\" → \"실험법\", \"실험 결과\" 등 매칭",
    "자동완성 제안 지원",
]
add_bullet_list(slide, Inches(0.9), Inches(2.1), Inches(5.3), Inches(2.5),
                search_items, font_size=13, spacing=Pt(7))

# 즐겨찾기
add_text_box(slide, Inches(7.0), Inches(1.65), Inches(5.5), Inches(0.35),
             "검색 기록 및 즐겨찾기", font_size=16, color=C_PRIMARY, bold=True)

fav_items = [
    "검색 기록:",
    "  • 검색어가 비어 있을 때 최근 검색 표시",
    "  • 클릭하면 해당 검색어로 재검색",
    "  • \"기록 삭제\" 버튼으로 전체 삭제",
    "",
    "검색어 즐겨찾기:",
    "  • 별표(★) 버튼으로 현재 검색어 저장",
    "  • 검색창 비어 있을 때 즐겨찾기 섹션 표시",
    "  • X 버튼으로 개별 삭제",
    "",
    "문서 즐겨찾기:",
    "  • 노트, 템플릿, 인벤토리 문서를 즐겨찾기 추가 가능",
]
add_bullet_list(slide, Inches(7.2), Inches(2.1), Inches(5.3), Inches(3.5),
                [i for i in fav_items if i], font_size=12, spacing=Pt(4))

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 14: 전자서명 1 — 개요
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "10", "전자서명 및 규정 준수", "Digital Signature & Compliance")
p = next_page()

# 서명 현황
add_text_box(slide, Inches(0.7), Inches(1.1), Inches(5.5), Inches(0.35),
             "서명 현황 페이지 (/signatures)", font_size=16, color=C_PRIMARY, bold=True)

# summary cards
sig_cards = [("서명 완료", "8", C_SUCCESS), ("서명 대기", "4", C_WARNING), ("잠긴 노트", "2", C_DANGER)]
for i, (label, val, color) in enumerate(sig_cards):
    x = Inches(0.7) + Inches(i * 1.7)
    add_rect(slide, x, Inches(1.6), Inches(1.5), Inches(0.8), fill_color=C_WHITE, line_color=color)
    add_text_box(slide, x + Inches(0.1), Inches(1.65), Inches(1.3), Inches(0.25),
                 label, font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(1.92), Inches(1.3), Inches(0.35),
                 val, font_size=20, color=color, bold=True, align=PP_ALIGN.CENTER)

add_table_slide(slide, Inches(0.7), Inches(2.65), [
    ["노트 상태", "버튼", "동작"],
    ["서명 완료 / 잠김", "\"열람\"", "읽기 전용으로 노트 보기"],
    ["진행 중", "\"서명\"", "노트 편집기로 이동 (서명 가능)"],
    ["초안", "\"수정\"", "노트 편집기로 이동"],
], [2.0, 1.2, 2.8], font_size=11)

# 서명 프로세스
add_text_box(slide, Inches(7.0), Inches(1.1), Inches(5.5), Inches(0.35),
             "전자서명 프로세스", font_size=16, color=C_PRIMARY, bold=True)

sign_steps = [
    "1. 노트 상태 확인: 반드시 \"진행 중(in_progress)\"",
    "2. 역할 확인: Reviewer 또는 Admin만 서명 가능",
    "3. 노트 편집기에서 \"서명\" 버튼 클릭",
    "4. 비밀번호 입력 대화상자에서 본인 비밀번호 입력",
    "5. 시스템이 자동 수행:",
    "   • 비밀번호 검증 (auth-service)",
    "   • 노트 내용 SHA-256 해시 계산",
    "   • 이전 서명 해시와 체이닝 (해시체인)",
    "   • 서명 레코드 저장",
    "   • 노트 상태 → 서명 완료(signed) 자동 전환",
    "6. 서명 완료 → 영구 불변 (수정/삭제/상태변경 불가)",
]
add_bullet_list(slide, Inches(7.2), Inches(1.6), Inches(5.3), Inches(4.0),
                sign_steps, font_size=12, spacing=Pt(5))

# 제약사항 박스
add_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5), fill_color=RGBColor(0xFE,0xF2,0xF2))
add_text_box(slide, Inches(0.7), Inches(5.1), Inches(4), Inches(0.3),
             "핵심 제약사항", font_size=14, color=C_DANGER, bold=True)
constraints = [
    "• Researcher는 서명 불가 (note:sign 권한 없음)",
    "• PATCH /status로 직접 \"signed\" 전환 불가 — 반드시 서명 프로세스 필요",
    "• 서명 완료(signed) 상태는 영구 불변 — 어떤 역할도 되돌릴 수 없음",
    "• DB 트리거 check_note_status_transition()이 잘못된 전환을 최종 차단",
]
add_bullet_list(slide, Inches(0.9), Inches(5.45), Inches(11.5), Inches(1.0),
                constraints, font_size=11, spacing=Pt(3), bullet_char="")

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 15: 감사 로그
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "11", "감사 로그", "Audit Logs")
p = next_page()

add_text_box(slide, Inches(0.7), Inches(1.1), Inches(11), Inches(0.35),
             "모든 시스템 활동이 자동으로 기록됩니다. 활동 기록은 수정/삭제할 수 없습니다.",
             font_size=13, color=C_GRAY)

# 필터
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.35),
             "필터 기능", font_size=16, color=C_PRIMARY, bold=True)
filter_items = [
    "수행자(actor) ID로 검색",
    "액션 유형 드롭다운 필터:",
    "  login, logout, note_created, note_updated,",
    "  note_status_changed, note_locked, signed,",
    "  note_admin_unlocked, export_requested 등",
    "날짜 범위 필터: 시작일 ~ 종료일",
    "페이지당 20건, 이전/다음 페이지 이동",
]
add_bullet_list(slide, Inches(0.9), Inches(2.05), Inches(5.3), Inches(2.5),
                filter_items, font_size=12, spacing=Pt(5))

# 표시 정보
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.35),
             "표시 정보", font_size=16, color=C_PRIMARY, bold=True)
add_table_slide(slide, Inches(7.0), Inches(2.05), [
    ["열", "설명"],
    ["시간", "발생 시각"],
    ["액션", "색상 뱃지로 구분"],
    ["수행자", "사용자 ID"],
    ["대상 유형/ID", "영향받은 엔티티"],
    ["상세", "key:value 형태의 부가 정보"],
    ["IP", "요청 IP 주소"],
], [1.5, 4.5], font_size=11)

# 기록되는 활동
add_text_box(slide, Inches(0.7), Inches(4.8), Inches(11), Inches(0.35),
             "서비스별 기록 액션", font_size=14, color=C_DARK, bold=True)
add_table_slide(slide, Inches(0.7), Inches(5.2), [
    ["서비스", "기록되는 액션"],
    ["auth-service", "로그인, 로그아웃, 사용자 생성/수정/삭제, 비밀번호 변경, 역할 변경"],
    ["eln-service", "노트 생성, 수정, 상태 변경, 잠금 해제, 삭제"],
    ["sig-audit-service", "전자서명, 서명 철회, 내보내기 요청"],
], [2.0, 10.0], font_size=11)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 16: 내보내기
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "12", "내보내기 (PDF / ZIP)", "Export")
p = next_page()

add_text_box(slide, Inches(0.7), Inches(1.1), Inches(11), Inches(0.35),
             "서명 완료(signed) 또는 잠김(locked) 상태의 노트를 PDF, ZIP 형식으로 내보낼 수 있습니다.",
             font_size=13, color=C_GRAY)

# 좌: 기능
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.35),
             "내보내기 기능", font_size=16, color=C_PRIMARY, bold=True)

export_items = [
    "전체 ZIP 내보내기:",
    "  모든 서명/잠금 노트를 하나의 ZIP으로 일괄 내보내기",
    "",
    "개별 노트 내보내기:",
    "  각 노트마다 PDF 또는 ZIP 버튼 제공",
    "",
    "진행 상태 표시:",
    "  요청 후 \"진행 중\" 섹션에 스피너 표시",
    "  3초마다 자동 상태 폴링",
    "  완료 시 \"다운로드\" 링크 / 실패 시 오류 메시지",
]
add_bullet_list(slide, Inches(0.9), Inches(2.05), Inches(5.3), Inches(3.5),
                [i for i in export_items if i], font_size=12, spacing=Pt(5))

# 우: 처리 과정
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.35),
             "내부 처리 과정", font_size=16, color=C_PRIMARY, bold=True)

# flow diagram as steps
flow_steps = [
    ("1", "요청 접수", "BullMQ 큐에 작업 등록"),
    ("2", "노트 조회", "워커가 eln-service에서 노트 내용 가져옴"),
    ("3", "PDF 렌더링", "Puppeteer로 HTML → PDF 변환"),
    ("4", "파일 업로드", "결과 파일을 MinIO에 저장"),
    ("5", "실시간 알림", "SSE로 클라이언트에 완료 알림 전달"),
]

for i, (num, title, desc) in enumerate(flow_steps):
    y = Inches(2.1) + Inches(i * 0.75)
    add_rect(slide, Inches(7.0), y, Inches(0.45), Inches(0.45), fill_color=C_PRIMARY)
    add_text_box(slide, Inches(7.0), y + Inches(0.07), Inches(0.45), Inches(0.3),
                 num, font_size=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.6), y + Inches(0.0), Inches(2.0), Inches(0.3),
                 title, font_size=13, color=C_DARK, bold=True)
    add_text_box(slide, Inches(7.6), y + Inches(0.28), Inches(4.5), Inches(0.3),
                 desc, font_size=11, color=C_GRAY)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 17: 알림
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "13", "알림", "Notifications")
p = next_page()

# 알림 유형
add_text_box(slide, Inches(0.7), Inches(1.1), Inches(5.5), Inches(0.35),
             "알림 유형", font_size=16, color=C_PRIMARY, bold=True)

add_table_slide(slide, Inches(0.7), Inches(1.55), [
    ["유형", "발생 시점", "수신자"],
    ["NOTE_LOCKED", "노트 잠금 시", "관련 사용자"],
    ["NOTE_SIGNED", "전자서명 완료 시", "관련 사용자"],
    ["NOTE_UNLOCKED", "관리자 잠금 해제 시", "관련 사용자"],
    ["BOOKING_APPROVED", "예약 승인 시", "예약 요청자"],
], [2.5, 2.5, 1.8], font_size=11)

# 알림 기능
add_text_box(slide, Inches(7.0), Inches(1.1), Inches(5.5), Inches(0.35),
             "알림 기능", font_size=16, color=C_PRIMARY, bold=True)

notif_items = [
    "헤더 우측 종(🔔) 아이콘 클릭 → 알림 패널",
    "미읽은 알림 수: 종 아이콘 위 빨간 뱃지",
    "알림 목록: 유형, 제목, 내용, 시간",
    "읽음 처리: 개별 또는 전체 읽음",
    "노트 알림 클릭 → 해당 노트로 이동",
    "예약 알림 클릭 → 예약 페이지로 이동",
    "30초마다 폴링 + SSE 실시간 수신",
]
add_bullet_list(slide, Inches(7.2), Inches(1.6), Inches(5.3), Inches(3.0),
                notif_items, font_size=13, spacing=Pt(7))

# 알림 전달 과정
add_text_box(slide, Inches(0.7), Inches(4.2), Inches(11), Inches(0.35),
             "알림 전달 과정", font_size=14, color=C_DARK, bold=True)

add_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.8), fill_color=C_BG_BLUE)
flow_text = (
    "이벤트 발생 (잠금/서명/해제/승인)\n"
    "  → 서비스에서 POST /api/notifications/internal 호출 (sig-audit-service)\n"
    "  → DB에 Notification 저장 + Redis Pub/Sub로 notifications:{userId} 채널에 발행\n"
    "  → API Gateway SSE 엔드포인트 (/api/events/notifications)가 구독\n"
    "  → 클라이언트 브라우저에 실시간 전달"
)
add_text_box(slide, Inches(0.7), Inches(4.85), Inches(11.9), Inches(1.5),
             flow_text, font_size=12, color=C_DARK)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 18: 관리자 기능
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "14", "관리자 기능", "Administration")
p = next_page()

add_text_box(slide, Inches(0.7), Inches(1.1), Inches(11), Inches(0.35),
             "Admin 역할의 사용자만 접근할 수 있는 관리 기능입니다.",
             font_size=13, color=C_GRAY)

# 4 cards
admin_cards = [
    ("사용자 관리\n/admin/users", [
        "사용자: 목록, 추가, 수정, 삭제",
        "  (이름, 이메일, 비밀번호, 역할, 상태)",
        "팀: 생성, 수정, 삭제",
        "  팀원 추가/제거 (검색 후 선택)",
        "조직: 생성, 수정, 삭제",
        "  (이름, 슬러그)",
    ], C_PRIMARY),
    ("역할 관리\n/admin/roles", [
        "역할 목록: 권한 뱃지 + 사용자 수",
        "역할 생성: 이름 + 권한 체크박스",
        "권한 수정: 드롭다운 → 체크박스 그리드",
        "역할 삭제: 확인 후 삭제",
        "4가지 기본 역할:",
        "  admin, researcher, reviewer, viewer",
    ], C_SUCCESS),
    ("코드 관리\n/admin/codes", [
        "참조 코드 그룹:",
        "  TEMPLATE_CATEGORY (프로토콜 카테고리)",
        "  INVENTORY_ITEM_TYPE (인벤토리 유형)",
        "코드 추가: 값, 표시명, 정렬 순서",
        "코드 수정: 활성/비활성 상태 전환",
        "코드 삭제: 확인 후 삭제",
    ], C_WARNING),
    ("시스템 설정\n/admin/settings", [
        "조직명: 현재 조직 이름 (읽기 전용)",
        "기본 언어: 한국어/영어 전환",
        "파일 저장소: MinIO 정보 (읽기 전용)",
        "",
        "",
        "",
    ], C_ACCENT),
]

for i, (title, items, color) in enumerate(admin_cards):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(col * 6.3)
    y = Inches(1.65) + Inches(row * 2.8)
    w = Inches(5.9)
    h = Inches(2.5)

    add_rect(slide, x, y, w, Inches(0.08), fill_color=color)
    add_rect(slide, x, y + Inches(0.08), w, h - Inches(0.08), fill_color=C_WHITE, line_color=RGBColor(0xE2,0xE8,0xF0))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), w - Inches(0.4), Inches(0.5),
                 title, font_size=13, color=color, bold=True)
    add_bullet_list(slide, x + Inches(0.2), y + Inches(0.7), w - Inches(0.4), h - Inches(0.8),
                    [it for it in items if it], font_size=11, color=C_DARK, spacing=Pt(4))

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 19: 역할 및 권한
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "14", "역할 및 권한 체계", "Roles & Permissions")
p = next_page()

# 역할 설명
add_text_box(slide, Inches(0.7), Inches(1.1), Inches(5.5), Inches(0.35),
             "기본 역할", font_size=16, color=C_PRIMARY, bold=True)
add_table_slide(slide, Inches(0.7), Inches(1.55), [
    ["역할", "설명"],
    ["Admin", "시스템 전체 관리자. 모든 기능 접근 가능"],
    ["Reviewer", "검토자. 노트 검토, 서명, 잠금 가능"],
    ["Researcher", "연구원. 노트 작성/수정 가능. 서명 불가"],
    ["Viewer", "열람자. 읽기 전용 접근"],
], [1.5, 5.0], font_size=12)

# 권한 비교표
add_text_box(slide, Inches(0.7), Inches(3.8), Inches(11), Inches(0.35),
             "역할별 주요 권한 비교", font_size=16, color=C_PRIMARY, bold=True)

add_table_slide(slide, Inches(0.7), Inches(4.25), [
    ["기능", "Admin", "Reviewer", "Researcher", "Viewer"],
    ["노트 조회", "O", "O", "O", "O"],
    ["노트 작성/수정", "O", "O", "O", "X"],
    ["노트 삭제", "O", "O", "O", "X"],
    ["전자서명", "O", "O", "X", "X"],
    ["in_progress → locked", "O", "O", "X", "X"],
    ["locked → draft (해제)", "O", "X", "X", "X"],
    ["관리자 메뉴", "O", "X", "X", "X"],
    ["예약 승인/반려", "O", "O", "X", "X"],
], [2.8, 1.0, 1.0, 1.2, 1.0], font_size=11)

# 권한 목록 (우측)
add_text_box(slide, Inches(7.5), Inches(1.1), Inches(5), Inches(0.35),
             "권한 목록 (18개)", font_size=16, color=C_PRIMARY, bold=True)
perms = [
    "note:read / write / delete / status / sign / unlock",
    "template:read / write",
    "inventory:read / write / delete",
    "scheduler:read / write",
    "audit:read",
    "export:pdf",
    "file:read / upload / delete",
]
add_bullet_list(slide, Inches(7.7), Inches(1.6), Inches(5), Inches(2.0),
                perms, font_size=12, spacing=Pt(5))

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 20: 노트 상태 흐름 상세
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "", "노트 상태 흐름 상세", "Note Status Lifecycle")
p = next_page()

# 상태 박스 그리기
states = [
    ("draft\n(초안)", Inches(1.5), Inches(2.5), RGBColor(0x64,0x74,0x8B)),
    ("in_progress\n(진행 중)", Inches(5.0), Inches(2.5), C_SECONDARY),
    ("locked\n(잠김)", Inches(9.0), Inches(1.8), C_DANGER),
    ("signed\n(서명 완료)", Inches(9.0), Inches(3.8), C_SUCCESS),
]

for name, x, y, color in states:
    add_rect(slide, x, y, Inches(2.0), Inches(1.0), fill_color=color)
    add_text_box(slide, x, y + Inches(0.15), Inches(2.0), Inches(0.7),
                 name, font_size=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# 화살표 설명
arrows = [
    (Inches(0.7), Inches(5.3), "draft → in_progress", "모든 역할", "PATCH /status"),
    (Inches(0.7), Inches(5.65), "in_progress → draft", "모든 역할", "PATCH /status"),
    (Inches(4.5), Inches(5.3), "in_progress → locked", "Reviewer, Admin", "PATCH /status"),
    (Inches(4.5), Inches(5.65), "in_progress → signed", "Reviewer, Admin", "POST /signatures/sign/:noteId"),
    (Inches(8.5), Inches(5.3), "locked → draft", "Admin (비밀번호)", "POST /admin-unlock"),
    (Inches(8.5), Inches(5.65), "signed → (모든 상태)", "불가", "—"),
]

add_table_slide(slide, Inches(0.7), Inches(5.0), [
    ["전환", "필요 역할", "API"],
] + [[a[2], a[3], a[4]] for a in arrows], [3.5, 2.5, 3.5], font_size=11)

# 상태별 제약 박스
add_rect(slide, Inches(0.5), Inches(7.6), Inches(12.3), Inches(0.0), fill_color=C_WHITE)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 21: FAQ
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
slide_title_bar(slide, "", "자주 묻는 질문 (FAQ)", "Frequently Asked Questions")
p = next_page()

faqs = [
    ("서명 완료된 노트를 수정할 수 있나요?",
     "아니요. 서명 완료(signed) 상태는 영구 불변입니다. 수정, 삭제, 상태 변경 모두 불가합니다."),
    ("잠긴 노트를 수정하려면?",
     "Admin이 잠금 해제 수행: 노트 메뉴 → \"잠금 해제\" → Admin 비밀번호 입력 → 초안(draft) 상태로 복귀"),
    ("Researcher로 서명할 수 있나요?",
     "아니요. Researcher에는 note:sign 권한이 없어 서명 요청 자체가 불가합니다."),
    ("예약 충돌은 어떻게 처리되나요?",
     "같은 시간대에 대기/승인 예약이 있으면 충돌 오류 발생. 승인 시 DB 비관적 락으로 동시 승인 방지."),
    ("인벤토리 알림은 어떻게 작동하나요?",
     "재고 부족: 수량 ≤ 최소수량 시 경고  |  만료 임박: 유효기간이 경고 일수(기본 30일) 이내 시 경고"),
    ("내보내기 파일은 어디에 저장되나요?",
     "MinIO 오브젝트 스토리지에 저장. 다운로드 링크는 시간 제한 있는 Presigned URL입니다."),
    ("검색이 잘 안 됩니다.",
     "인덱스 갱신은 비동기 처리되어 생성/수정 직후 약간의 지연이 있을 수 있습니다."),
    ("다국어 전환은?",
     "헤더 우측 언어(🌐) 버튼 클릭 → 한국어/영어 전환. 설정은 브라우저에 저장됩니다."),
]

for i, (q, a) in enumerate(faqs):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(col * 6.3)
    y = Inches(1.1) + Inches(row * 1.45)
    w = Inches(5.9)

    add_rect(slide, x, y, w, Inches(1.3), fill_color=C_WHITE, line_color=RGBColor(0xE2,0xE8,0xF0))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.08), w - Inches(0.3), Inches(0.35),
                 f"Q. {q}", font_size=11, color=C_PRIMARY, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.45), w - Inches(0.3), Inches(0.8),
                 a, font_size=10, color=C_DARK)

page_number(slide, p)

# ════════════════════════════════════════════════════════════════
# SLIDE 22: 마지막 — 감사
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_rect(slide, 0, 0, W, H, C_PRIMARY, RGBColor(0x0E, 0x2A, 0x6B))

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(0.8),
             "감사합니다", font_size=44, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.5),
             "LabNote ELN 사용자 매뉴얼 v1.0", font_size=20, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.3), Inches(2.3), Pt(2))
line2.fill.solid()
line2.fill.fore_color.rgb = C_ACCENT
line2.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(4.7), Inches(10), Inches(0.4),
             "문의사항은 시스템 관리자에게 연락해 주세요.", font_size=14, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

# ── 저장 ──
output_path = r"C:\vscode\lab\lab-companion-feature-phase1-backend-infra\services\LabNote_ELN_사용자매뉴얼.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
